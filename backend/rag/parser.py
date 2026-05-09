import os
from typing import List
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain.schema import Document


def parse_pdf(file_path: str) -> List[Document]:
    from pypdf import PdfReader
    reader = PdfReader(file_path)
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
    return [Document(page_content=text, metadata={"source": os.path.basename(file_path), "type": "pdf"})]


def parse_docx(file_path: str) -> List[Document]:
    from docx import Document as DocxDocument
    doc = DocxDocument(file_path)
    text = "\n".join([para.text for para in doc.paragraphs])
    return [Document(page_content=text, metadata={"source": os.path.basename(file_path), "type": "docx"})]


def parse_pptx(file_path: str) -> List[Document]:
    from pptx import Presentation
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return [Document(page_content=text, metadata={"source": os.path.basename(file_path), "type": "pptx"})]


def parse_text(file_path: str) -> List[Document]:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    return [Document(page_content=text, metadata={"source": os.path.basename(file_path), "type": "text"})]


PARSERS = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".pptx": parse_pptx,
    ".txt": parse_text,
    ".md": parse_text,
    ".py": parse_text,
    ".js": parse_text,
    ".ts": parse_text,
    ".java": parse_text,
    ".c": parse_text,
    ".cpp": parse_text,
    ".h": parse_text,
    ".csv": parse_text,
    ".json": parse_text,
}


def parse_file(file_path: str) -> List[Document]:
    ext = os.path.splitext(file_path)[1].lower()
    parser = PARSERS.get(ext)
    if not parser:
        raise ValueError(f"Unsupported file type: {ext}")
    return parser(file_path)


def chunk_documents(documents: List[Document], chunk_size: int = 1000, chunk_overlap: int = 200) -> List[Document]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        length_function=len,
    )
    return splitter.split_documents(documents)
